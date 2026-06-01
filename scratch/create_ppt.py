# -*- coding: utf-8 -*-
"""
Generate an ultra-premium, sleek Cyber-Dark Tech PowerPoint presentation
for institutional-level Science & Technology research proposal (Đề tài KHCN cấp Cơ sở).
Fixes:
  - Added a brand new dedicated Slide 3: "CÁC NGHIÊN CỨU TIỀN ĐỀ LIÊN QUAN" (Related Work)
    to summarize the academic foundation of the research proposal.
  - Formatted using the clean 3-column card grid, no bullet points, and high legibility.
  - Saved to the standard file: `de_cuong_bao_cao_chuan.pptx`.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Sleek Cyber-Dark Color Palette ───────────────────────────────
BG           = RGBColor(10, 15, 30)     # Deep Cyber Space Dark Background
DARK_CARD    = RGBColor(20, 27, 45)     # Lighter Dark Indigo for Glassmorphic Cards
CYAN         = RGBColor(6, 182, 212)    # Glowing Cyan (Primary Neon)
EMERALD      = RGBColor(16, 185, 129)   # Glowing Emerald (Secondary Accent)
WHITE        = RGBColor(248, 250, 252)  # High-contrast clean white text
GRAY_MUTED   = RGBColor(148, 163, 184)  # Muted slate gray for secondary descriptions
WARN_ROSE    = RGBColor(244, 63, 94)    # Warning Rose Red
WARN_CARD    = RGBColor(34, 25, 38)     # Soft dark red card fill

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_FAMILY = "Segoe UI"
BASE_PATH = r"e:\HK2(2025-2026)\Đồ Án 2\MainFolder\MyProject"

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_top_bar(slide, height=Inches(0.08)):
    """Thin elegant glowing top bar."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, height)  # 1 = Rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()

def add_bottom_bar(slide, text="", height=Inches(0.4)):
    """Sleek minimalist footer."""
    y = SLIDE_H - height
    bar = slide.shapes.add_shape(1, 0, y, SLIDE_W, height)  # 1 = Rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 20, 38)
    bar.line.fill.background()
    if text:
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.color.rgb = GRAY_MUTED
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_title(slide, text, top=Inches(0.35), left=Inches(0.8), width=Inches(11.7), size=Pt(26)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_FAMILY
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_icon_badge(slide, emoji, top, left, bg_color=DARK_CARD, border_color=CYAN, size=Inches(0.45)):
    """Draws a beautiful circular sticker/badge with a high-fidelity tech emoji inside (Shape ID 9 = Oval)."""
    circle = slide.shapes.add_shape(9, left, top, size, size)  # 9 = Oval (Circle when width == height)
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg_color
    circle.line.color.rgb = border_color
    circle.line.width = Pt(1.5)
    
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.name = "Segoe UI Symbol"
    p.font.size = Pt(15)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def draw_premium_card(slide, title, paragraphs, top, left, width, height, bg_color=DARK_CARD, glow_color=CYAN, title_size=16, body_size=14.0, emoji=None):
    """
    Draws a premium rounded card with a glowing border outline following the card's curves perfectly.
    """
    # 1. Main Card Body
    card = slide.shapes.add_shape(5, left, top, width, height)  # 5 = Rounded Rectangle
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = glow_color
    card.line.width = Pt(1.5)  # Glowing outline border
    
    # 2. Add Emoji Badge if provided
    text_shift = Inches(0.0)
    if emoji:
        add_icon_badge(slide, emoji, top + Inches(0.18), left + Inches(0.25), bg_color=bg_color, border_color=glow_color)
        text_shift = Inches(0.55)
    
    # 3. Content Text Box
    tb = slide.shapes.add_textbox(left + Inches(0.2) + text_shift, top + Inches(0.15), width - Inches(0.35) - text_shift, height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    start_idx = 0
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = glow_color
        p.space_after = Pt(8)
        start_idx = 1
        
    for i, item in enumerate(paragraphs):
        p = tf.add_paragraph() if (i > 0 or start_idx == 1) else tf.paragraphs[0]
        
        if isinstance(item, str):
            p.text = item
            p.font.size = Pt(body_size)
            p.font.color.rgb = WHITE
            p.font.bold = False
        elif isinstance(item, dict):
            p.text = item.get("text", "")
            p.font.size = Pt(item.get("size", body_size))
            p.font.bold = item.get("bold", False)
            p.font.italic = item.get("italic", False)
            p.font.color.rgb = item.get("color", WHITE)
            p.space_before = Pt(item.get("space_before", 6))
            p.space_after = Pt(item.get("space_after", 0))
            
        p.font.name = FONT_FAMILY
        p.line_spacing = 1.15

def embed_diagram(slide, filename, top, left, width, height=None):
    """Programmatically embeds actual diagram image files directly into the slide layouts."""
    img_path = os.path.join(BASE_PATH, filename)
    if os.path.exists(img_path):
        if height:
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top, width=width)
    else:
        box = slide.shapes.add_shape(5, left, top, width, height or Inches(5.4))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_CARD
        box.line.color.rgb = WARN_ROSE
        box.line.width = Pt(1.5)
        
        tb = slide.shapes.add_textbox(left, top + Inches(2.0), width, Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ WARNING ]\nMissing {filename} in workspace!"
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.color.rgb = WARN_ROSE
        p.alignment = PP_ALIGN.CENTER

FOOTER = "Đề tài KHCN cấp Cơ sở  |  Ứng dụng LLM xây dựng Đồ thị Tri thức hỗ trợ chẩn đoán bệnh Tiểu đường  |  Trường Đại học Nam Cần Thơ"

# ══════════════════════════════════════════════════════════════════
prs = new_prs()

# ── SLIDE 1: Split Cyber-Dark Cover ────────────────────────────────
s = add_slide(prs)
add_top_bar(s, Inches(0.1))
add_bottom_bar(s, "Trường Công Nghệ Số Và Trí Tuệ Nhân Tạo  ·  Trường Đại Học Nam Cần Thơ  ·  2026")

left_decor = s.shapes.add_shape(1, 0, 0, Inches(0.15), SLIDE_H)
left_decor.fill.solid()
left_decor.fill.fore_color.rgb = CYAN
left_decor.line.fill.background()

tb_title = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.5))
tf_title = tb_title.text_frame
tf_title.word_wrap = True

p_sub = tf_title.paragraphs[0]
p_sub.text = "ĐỀ CƯƠNG ĐĂNG KÝ ĐỀ TÀI KH&CN CẤP CƠ SỞ"
p_sub.font.name = FONT_FAMILY
p_sub.font.size = Pt(15)
p_sub.font.bold = True
p_sub.font.color.rgb = CYAN
p_sub.space_after = Pt(15)

p_main = tf_title.add_paragraph()
p_main.text = "Ứng Dụng LLM Xây Dựng Đồ Thị Tri Thức\nHỗ Trợ Chẩn Đoán Bệnh Tiểu Đường"
p_main.font.name = FONT_FAMILY
p_main.font.size = Pt(38)
p_main.font.bold = True
p_main.font.color.rgb = WHITE
p_main.space_after = Pt(40)

p_team_lead = tf_title.add_paragraph()
p_team_lead.text = "Chủ nhiệm đề tài: Lê Quang Huy (Lớp: DH22KPM01  ·  Khóa: 10)"
p_team_lead.font.name = FONT_FAMILY
p_team_lead.font.size = Pt(16)
p_team_lead.font.bold = True
p_team_lead.font.color.rgb = EMERALD
p_team_lead.space_after = Pt(8)

p_members = tf_title.add_paragraph()
p_members.text = "Thành viên nghiên cứu: Lưu An Thuận (Thư ký)  ·  Lưu Vỹ Thuận  ·  Nguyễn Văn Hồ"
p_members.font.name = FONT_FAMILY
p_members.font.size = Pt(15)
p_members.font.color.rgb = WHITE
p_members.space_after = Pt(8)

p_adv = tf_title.add_paragraph()
p_adv.text = "Giảng viên hướng dẫn: ThS. Đặng Mạnh Huy"
p_adv.font.name = FONT_FAMILY
p_adv.font.size = Pt(15)
p_adv.font.bold = True
p_adv.font.color.rgb = CYAN


# ── SLIDE 2: 1. TỔNG QUAN & ĐẶT VẤN ĐỀ ─────────────────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "1. TỔNG QUAN & ĐẶT VẤN ĐỀ")

draw_premium_card(
    s,
    "Bối cảnh thực tiễn",
    [
        "Khai thác y văn tự động đòi hỏi độ chính xác tuyệt đối.",
        "Đồ thị tri thức (KG) cung cấp biểu diễn tri thức y sinh minh bạch, làm nền tảng cho lập luận lâm sàng giải thích được (Explainable AI)."
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(5.8), height=Inches(2.4),
    glow_color=CYAN, emoji="🧬"
)

draw_premium_card(
    s,
    "Khoảng trống nghiên cứu",
    [
        "Chưa có quy trình tự động hóa dựng Đồ thị tri thức chuyên sâu đái tháo đường theo framework EDC ràng buộc siêu từ điển UMLS.",
        "Thiếu vắng sự phối hợp với cơ chế phản biện đa tác nhân đồng cấp (P2P Debate Gate) tích hợp luật Veto lâm sàng để triệt tiêu ảo giác."
    ],
    top=Inches(3.8), left=Inches(0.6), width=Inches(5.8), height=Inches(2.8),
    glow_color=EMERALD, emoji="🎯"
)

draw_premium_card(
    s,
    "Rào cản chí mạng khi dùng LLM độc lập",
    [
        {"text": "Hiện tượng ảo giác tri thức (hallucination):", "bold": True, "color": WARN_ROSE, "space_before": 4},
        {"text": "LLM dễ sinh câu trả lời tự tin nhưng không chính xác với y văn lâm sàng thực tế [1].", "space_before": 4},
        {"text": "Thiếu minh chứng lâm sàng:", "bold": True, "color": WARN_ROSE, "space_before": 12},
        {"text": "Không có khả năng đối khớp nguồn gốc của thực thể và quan hệ y sinh được trích xuất.", "space_before": 4},
        {"text": "Không có ràng buộc lược đồ (Schema):", "bold": True, "color": WARN_ROSE, "space_before": 12},
        {"text": "Trích xuất tự do dễ dẫn đến các bộ ba tri thức phi logic y sinh học.", "space_before": 4}
    ],
    top=Inches(1.2), left=Inches(6.8), width=Inches(5.9), height=Inches(5.4),
    bg_color=WARN_CARD, glow_color=WARN_ROSE, emoji="⚠️"
)


# ── SLIDE 3: 1b. CÁC NGHIÊN CỨU TIỀN ĐỀ LIÊN QUAN (New Slide) ──────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "1b. CÁC NGHIÊN CỨU TIỀN ĐỀ LIÊN QUAN")

# Column 1
draw_premium_card(
    s,
    "Kế thừa Lược đồ ràng buộc EDC",
    [
        {"text": "Thành tựu tiền đề:", "bold": True, "color": CYAN},
        {"text": "Xây dựng thành công quy trình trích xuất ràng buộc 3 bước (Extract-Define-Canonicalize) giúp tự động hóa dựng đồ thị tri thức y sinh [4][5].", "size": 11.0},
        {"text": "Hạn chế & Khoảng trống:", "bold": True, "color": CYAN, "space_before": 8},
        {"text": "Chưa tối ưu hóa cho chuyên khoa tiểu đường phức tạp; thiếu cơ chế chuẩn hóa sâu để khử trùng lặp thực thể về siêu từ điển y học quy mô lớn.", "size": 11.0},
        {"text": "Giải pháp cải tiến đề tài:", "bold": True, "color": CYAN, "space_before": 8},
        {"text": "Kế thừa quy trình EDC, phát triển module tự động hóa ánh xạ ngữ nghĩa chuẩn hóa toàn bộ thực thể thô về mã UMLS CUI duy nhất ở Phase 3b.", "size": 11.0}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(3.8), height=Inches(5.4),
    glow_color=CYAN, emoji="🧬"
)

# Column 2
draw_premium_card(
    s,
    "Học hỏi mô hình Local GraphRAG",
    [
        {"text": "Thành tựu tiền đề:", "bold": True, "color": EMERALD},
        {"text": "Ứng dụng thành công GraphRAG lâm sàng trên mô hình ngôn ngữ lớn cục bộ (Local LLM) để bảo mật thông tin y tế chuyên sâu [2][3].", "size": 11.0},
        {"text": "Hạn chế & Khoảng trống:", "bold": True, "color": EMERALD, "space_before": 8},
        {"text": "Mới chỉ dừng lại ở các tác vụ hỏi đáp cơ bản; chưa tích hợp hệ thống ngắt mạch an toàn lâm sàng thời gian thực khi có xung đột phác đồ hoặc suy chức năng tạng.", "size": 11.0},
        {"text": "Giải pháp cải tiến đề tài:", "bold": True, "color": EMERALD, "space_before": 8},
        {"text": "Mở rộng động cơ GraphRAG cục bộ trên y văn Merck Manuals và xây dựng bộ ngắt mạch y khoa (Clinical Circuit Breaker) tự động phát hiện rủi ro để bảo vệ bệnh nhân.", "size": 11.0}
    ],
    top=Inches(1.2), left=Inches(4.7), width=Inches(3.9), height=Inches(5.4),
    glow_color=EMERALD, emoji="⚙️"
)

# Column 3
draw_premium_card(
    s,
    "Áp dụng phản biện Multi-Agent",
    [
        {"text": "Thành tựu tiền đề:", "bold": True, "color": WARN_ROSE},
        {"text": "Đề xuất giao thức tranh luận đa tác nhân (Multi-Agent Debate) đồng cấp giúp giảm thiểu ảo giác tri thức của mô hình ngôn ngữ lớn [6][7][8].", "size": 11.0},
        {"text": "Hạn chế & Khoảng trống:", "bold": True, "color": WARN_ROSE, "space_before": 8},
        {"text": "Mới đánh giá các tác vụ NLP hội thoại tổng quát; chưa xây dựng luật ngắt mạch y sinh học hoặc công thức chốt điểm đồng thuận để lọc dữ liệu trước khi nạp đồ thị.", "size": 11.0},
        {"text": "Giải pháp cải tiến đề tài:", "bold": True, "color": WARN_ROSE, "space_before": 8},
        {"text": "Hiện thực hóa giao thức thành cổng tự động thẩm định P2P Debate Gate, chốt điểm đồng thuận FCS >= 80 và tích hợp Luật phủ quyết lâm sàng (Clinical Veto).", "size": 11.0}
    ],
    top=Inches(1.2), left=Inches(8.9), width=Inches(3.8), height=Inches(5.4),
    glow_color=WARN_ROSE, emoji="👥"
)


# ── SLIDE 4: 2. MỤC TIÊU & PHẠM VI NGHIÊN CỨU ──────────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "2. MỤC TIÊU & PHẠM VI NGHIÊN CỨU")

draw_premium_card(
    s,
    "Mục tiêu nghiên cứu",
    [
        {"text": "Mục tiêu chính:", "bold": True, "color": EMERALD},
        {"text": "Ứng dụng LLM xây dựng Đồ thị tri thức chuẩn hóa UMLS hỗ trợ chẩn đoán đái tháo đường; thiết lập quy trình thẩm định chất lượng tri thức tự động dựa trên giao thức Đa tác nhân (Multi-Agent)."},
        {"text": "Mục tiêu cụ thể:", "bold": True, "color": EMERALD, "space_before": 12},
        {"text": "1. Tiền xử lý tự động dữ liệu y văn từ cẩm nang Merck Manuals Professional Version.", "space_before": 6},
        {"text": "2. Trích xuất tri thức theo EDC framework, chuẩn hóa UMLS CUI.", "space_before": 4},
        {"text": "3. Thiết lập & tối ưu hóa cơ sở dữ liệu đồ thị Neo4j.", "space_before": 4},
        {"text": "4. Đánh giá chất lượng đồ thị qua giao thức phản biện Multi-Agent.", "space_before": 4},
        {"text": "5. Triển khai Web demo CDSS tích hợp ngắt mạch y khoa an toàn.", "space_before": 4}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(5.9), height=Inches(5.4),
    glow_color=EMERALD, emoji="🚀"
)

draw_premium_card(
    s,
    "Phạm vi & Giới hạn thực nghiệm",
    [
        {"text": "Đối tượng nghiên cứu:", "bold": True, "color": CYAN},
        {"text": "Quy trình trích xuất, chuẩn hóa tri thức (UMLS CUI) và giao thức tự động hóa thẩm định chất lượng đồ thị bằng Đa tác nhân chuyên ngành đái tháo đường.", "space_before": 6},
        {"text": "Dữ liệu thực nghiệm:", "bold": True, "color": CYAN, "space_before": 14},
        {"text": "Y văn điều trị đái tháo đường từ Cẩm nang chuyên môn Merck Manuals Professional Version.", "space_before": 6},
        {"text": "Giới hạn thực nghiệm:", "bold": True, "color": WARN_ROSE, "space_before": 14},
        {"text": "Giao diện CDSS và tính năng ngắt mạch tương tác thuốc chỉ phát triển ở quy mô demo thực nghiệm giả lập trong phòng Lab để chứng minh tính ứng dụng thực tế.", "space_before": 6}
    ],
    top=Inches(1.2), left=Inches(6.8), width=Inches(5.9), height=Inches(5.4),
    glow_color=CYAN, emoji="🛡️"
)


# ── SLIDE 5: 3. PHƯƠNG PHÁP NGHIÊN CỨU ────────────────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "3. PHƯƠNG PHÁP NGHIÊN CỨU")

draw_premium_card(
    s,
    "Nghiên cứu lý thuyết",
    [
        {"text": "Thiết kế Graph Schema:", "bold": True, "color": CYAN},
        {"text": "Khảo sát cẩm nang y khoa Merck Manuals & siêu từ điển y học UMLS để lập lược đồ đồ thị chuẩn."},
        {"text": "Công nghệ y khoa AI:", "bold": True, "color": CYAN, "space_before": 12},
        {"text": "Nghiên cứu lý thuyết Đồ thị tri thức (KG), mô hình nhúng y sinh (BioBERT), kiến trúc GraphRAG và cơ chế đàm thoại phản biện đa tác nhân (Multi-Agent Debate)."}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(3.8), height=Inches(5.4),
    glow_color=CYAN, emoji="🧠", body_size=13.5
)

draw_premium_card(
    s,
    "Phương pháp thực nghiệm",
    [
        {"text": "Tiền xử lý & làm sạch:", "bold": True, "color": EMERALD},
        {"text": "Thu thập dữ liệu y văn Merck Manuals, tự động loại bỏ ký tự nhiễu; ứng dụng kỹ thuật đối khớp ngữ nghĩa để chuẩn hóa từ đồng nghĩa về mã UMLS CUI duy nhất."},
        {"text": "Hiện thực hóa hệ thống:", "bold": True, "color": EMERALD, "space_before": 12},
        {"text": "Cài đặt pipeline trích xuất EDC; nạp CSDL đồ thị Neo4j bằng Cypher; xây dựng ứng dụng web CDSS (FastAPI kết hợp ReactJS)."}
    ],
    top=Inches(1.2), left=Inches(4.7), width=Inches(3.9), height=Inches(5.4),
    glow_color=EMERALD, emoji="💻", body_size=13.5
)

draw_premium_card(
    s,
    "Đánh giá & Kiểm chứng",
    [
        {"text": "Thẩm định tự động:", "bold": True, "color": WARN_ROSE},
        {"text": "Phát triển hệ thống Hội đồng đa tác nhân đồng cấp (P2P Debate Gate) tự đàm thoại và chốt điểm đồng thuận chéo."},
        {"text": "Đo lường định lượng:", "bold": True, "color": WARN_ROSE, "space_before": 12},
        {"text": "Tính toán Precision, Recall, F1-Score trên 3 cấp độ (Exact, Strict, Partial). Chạy thử nghiệm giả lập ca lâm sàng thực tế để kiểm định độ nhạy của bộ ngắt mạch Circuit Breaker."}
    ],
    top=Inches(1.2), left=Inches(8.9), width=Inches(3.8), height=Inches(5.4),
    glow_color=WARN_ROSE, emoji="🧪", body_size=13.5
)


# ── SLIDE 6: 4. NỘI DUNG 1 & 2 ────────────────────────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "4. NỘI DUNG 1 & 2: LÝ THUYẾT + TIỀN XỬ LÝ DỮ LIỆU")

draw_premium_card(
    s,
    "Nội dung 1: Khảo sát cơ sở lý thuyết",
    [
        {"text": "Đồ thị tri thức (KG):", "bold": True, "color": CYAN},
        {"text": "Tìm hiểu cấu trúc mạng lưới liên kết gồm Thực thể (Bệnh lý, Thuốc, Chỉ số eGFR), Mối quan hệ (Chống chỉ định, Điều trị) và các Thuộc tính ngữ cảnh lâm sàng."},
        {"text": "Kiến trúc GraphRAG & Multi-Agent:", "bold": True, "color": CYAN, "space_before": 12},
        {"text": "Kết hợp CSDL đồ thị và LLM để giữ nguyên ngữ cảnh y văn; nghiên cứu cơ chế đàm thoại phản biện (Debate Paradigm) giữa các Agent để giảm thiểu tối đa hiện tượng ảo giác lâm sàng."}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(5.9), height=Inches(4.3),
    glow_color=CYAN, emoji="📖"
)

draw_premium_card(
    s,
    "Nội dung 2: Tiền xử lý dữ liệu lâm sàng y văn",
    [
        {"text": "Thu thập và phân đoạn văn bản:", "bold": True, "color": EMERALD},
        {"text": "Tập trung khai thác nguồn dữ liệu tri thức y khoa dạng văn xuôi (narrative prose) có cấu trúc ngữ nghĩa chặt chẽ từ Cẩm nang chuyên môn Merck Manuals chuyên mục Điều trị Đái tháo đường."},
        {"text": "Làm sạch ký tự và chuẩn hóa dữ liệu:", "bold": True, "color": EMERALD, "space_before": 12},
        {"text": "Triển khai module tiền xử lý văn bản tự động loại bỏ thẻ HTML, ký hiệu thừa khi trích xuất web, tối ưu hóa cấu trúc văn bản lâm sàng sạch cho LLM."}
    ],
    top=Inches(1.2), left=Inches(6.8), width=Inches(5.9), height=Inches(4.3),
    glow_color=EMERALD, emoji="📥"
)

# Bottom alert banner for Results
draw_premium_card(
    s,
    "",
    [{"text": "→  Ý nghĩa: Xây dựng dòng dữ liệu lâm sàng cực kỳ sạch sẽ, cấu trúc hóa cao, tạo tiền đề tối ưu cho việc trích xuất tri thức bằng mô hình ngôn ngữ lớn.", "bold": True, "italic": True, "color": EMERALD}],
    top=Inches(5.7), left=Inches(0.6), width=Inches(12.1), height=Inches(0.8),
    bg_color=RGBColor(16, 38, 32), glow_color=EMERALD
)


# ── SLIDE 7: 5. NỘI DUNG 3 – EDC PIPELINE ─────────────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "5. NỘI DUNG 3: TRÍCH XUẤT RÀNG BUỘC THEO EDC FRAMEWORK")

draw_premium_card(
    s,
    "Framework trích xuất có ràng buộc EDC [5]",
    [
        {"text": "Phase 1 - Extract (Clinical Open IE):", "bold": True, "color": CYAN, "space_before": 4},
        {"text": "Sử dụng Few-shot Prompt tối ưu để LLM trích xuất mở các bộ ba thô dạng (Bệnh lý - treats/contraindicates - Dược chất) từ hồ sơ và phác đồ."},
        {"text": "Phase 2 - Define (Semantic Definition):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Gọi LLM tự động viết định nghĩa rõ ràng cho từng thực thể và quan hệ dựa theo ngữ cảnh xuất hiện, giúp bảo toàn toàn vẹn ngữ nghĩa lâm sàng gốc."},
        {"text": "Phase 3 - Canonicalize (Chuẩn hóa Lược đồ):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "3a - Relation: Ánh xạ quan hệ thô về lược đồ chuẩn (diabetes_schema.csv) bằng Cosine Similarity.\n3b - Entity-Type: Ánh xạ kiểu thực thể về mã UMLS CUI chuẩn (CUI mapping)."},
        {"text": "Phase 4 - Property Packing (Đóng gói thuộc tính):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Module PropertyPacker truy vấn ngoại tuyến gán mã chuẩn y khoa RxNorm (thuốc), ICD-10 (bệnh lý), khử trùng lặp và đóng gói thành tệp JSON đồ thị chuẩn."}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(7.3), height=Inches(5.4),
    glow_color=CYAN, emoji="⚙️", body_size=13.5
)

# Programmatically embed the real EDC flow diagram image
embed_diagram(s, "edc_pipeline_actual_flow.png", top=Inches(1.2), left=Inches(8.1), width=Inches(4.6), height=Inches(5.4))


# ── SLIDE 8: 6. NỘI DUNG 4 – HỘI ĐỒNG ĐA TÁC NHÂN ──────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "6. NỘI DUNG 4: THẨM ĐỊNH QUA HỘI ĐỒNG ĐA TÁC NHÂN (P2P DEBATE GATE)")

draw_premium_card(
    s,
    "Kiến trúc hội đồng phản biện y khoa đa tác nhân [7][8]",
    [
        {"text": "1. Clinical Specialist Agent (Llama-3.3-70B - Trọng số 0.4):", "bold": True, "color": EMERALD, "space_before": 4},
        {"text": "Đánh giá tính chính xác sinh y học, an toàn lâm sàng, đối chiếu Cẩm nang điều trị đái tháo đường Merck Manuals."},
        {"text": "2. Ontology Inspector Agent (Llama-3.1-8B - Trọng số 0.3):", "bold": True, "color": EMERALD, "space_before": 10},
        {"text": "Kiểm tra tính tuân thủ nghiêm ngặt của các bộ ba tri thức với lược đồ Domain/Range định sẵn trong diabetes_schema.csv."},
        {"text": "3. Medical Skeptic Agent (Gemma-2-9B - Trọng số 0.3):", "bold": True, "color": EMERALD, "space_before": 10},
        {"text": "Đóng vai trò phản biện độc lập, tìm kiếm lỗi ngôn ngữ NLP, triệt tiêu ảo giác tri thức và các thực thể nhiễu."},
        {"text": "Cơ chế Điểm đồng thuận (FCS) & Luật Phủ quyết (Veto Rule):", "bold": True, "color": WARN_ROSE, "space_before": 14},
        {"text": "Điểm đồng thuận FCS = 0.4×CS + 0.3×OI + 0.3×MS. Triplet được lưu vào Neo4j nếu FCS ≥ 80/100.\n\nLuật Phủ quyết: Bất kỳ Agent nào đưa ra kết luận bác bỏ [INCORRECT] với độ tin cậy > 70% sẽ lập tức loại bỏ triplet khỏi luồng đóng gói."}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(7.3), height=Inches(5.4),
    glow_color=EMERALD, emoji="👥", body_size=13.5
)

# Programmatically embed the real Multi-Agent debate framework diagram image
embed_diagram(s, "multiagent_debate_framework.png", top=Inches(1.2), left=Inches(8.1), width=Inches(4.6), height=Inches(5.4))


# ── SLIDE 9: 7. NỘI DUNG 5 – CDSS & CIRCUIT BREAKER ───────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "7. NỘI DUNG 5: ỨNG DỤNG CDSS & CIRCUIT BREAKER (THỰC NGHIỆM)")

draw_premium_card(
    s,
    "Mô hình CDSS thực nghiệm dựa trên Neo4j GraphRAG",
    [
        {"text": "Backend API y khoa điều phối song song (FastAPI):", "bold": True, "color": CYAN, "space_before": 4},
        {"text": "Đồng thời thực hiện Vector Search y văn và truy vấn cấu trúc đồ thị Neo4j bằng ngôn ngữ Cypher, đảm bảo tính nhất quán của tri thức lâm sàng."},
        {"text": "Động cơ GraphRAG hỗ trợ chẩn đoán:", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Bóc tách các tuyến tri thức liên kết (Knowledge Paths) từ Neo4j nạp vào prompt để LLM sinh chẩn đoán chính xác nhất."},
        {"text": "Ngắt mạch an toàn lâm sàng (Clinical Circuit Breaker):", "bold": True, "color": WARN_ROSE, "space_before": 10},
        {"text": "Module tự động quét bệnh án, đối chiếu chống chỉ định và tương tác thuốc trong thời gian thực. Phát hiện rủi ro lập tức ngắt hành động, hiển thị Cảnh báo đỏ y khoa (Red Alert)."},
        {"text": "Giao diện Web tương tác (ReactJS):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Tích hợp màn hình chat y khoa, bản đồ trực quan hóa đồ thị Neo4j, bảng thông số lâm sàng bệnh nhân EHR và ô cảnh báo ngắt mạch."}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(7.3), height=Inches(5.4),
    glow_color=CYAN, emoji="🩺", body_size=13.5
)

# Programmatically embed the real End-to-End CDSS & KG pipeline diagram image
embed_diagram(s, "e2e_kg_pipeline.png", top=Inches(1.2), left=Inches(8.1), width=Inches(4.6), height=Inches(5.4))


# ── SLIDE 10: 8. PHÂN CÔNG & TIẾN ĐỘ ───────────────────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "8. PHÂN CÔNG NHIỆM VỤ & TIẾN ĐỘ THỰC HIỆN")

# Left Column: Roles
draw_premium_card(
    s,
    "Phân công nhiệm vụ thành viên",
    [
        {"text": "Lê Quang Huy (Chủ nhiệm đề tài):", "bold": True, "color": EMERALD, "space_before": 4},
        {"text": "Thiết kế tổng thể thuật toán EDC, GraphRAG; quản lý tích hợp; định hướng khoa học và viết báo cáo thuyết minh.", "space_before": 4},
        {"text": "Lưu An Thuận (Thư ký):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Khai thác tài liệu Merck Manuals; tiền xử lý làm sạch y văn; thiết lập và quản trị cơ sở dữ liệu Neo4j.", "space_before": 4},
        {"text": "Lưu Vỹ Thuận (Thành viên):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Thiết kế Prompt LLM; xây dựng module định nghĩa NLG; lập trình tích hợp Circuit Breaker lên web demo.", "space_before": 4},
        {"text": "Nguyễn Văn Hồ (Thành viên):", "bold": True, "color": CYAN, "space_before": 10},
        {"text": "Xây dựng hệ thống thẩm định đa tác nhân (Agent Debate); viết module đo lường Precision, Recall, F1-Score.", "space_before": 4}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(5.9), height=Inches(5.4),
    glow_color=EMERALD, emoji="👥", body_size=13.5
)

# Right Column: Timeline
draw_premium_card(
    s,
    "Tiến độ thực hiện đề tài (6 tháng)",
    [
        {"text": "Tháng 04/2026:", "bold": True, "color": CYAN, "space_before": 4},
        {"text": "Thu thập và phân đoạn y văn Merck Manuals.\nThiết lập quy trình tự động làm sạch ký tự y khoa nhiễu.\nĐịnh nghĩa và thiết kế Graph Schema chuẩn hóa.", "space_before": 4},
        {"text": "Tháng 05/2026:", "bold": True, "color": CYAN, "space_before": 14},
        {"text": "Cài đặt quy trình trích xuất ràng buộc EDC.\nXây dựng Hội đồng đa tác nhân thẩm định đồng thuận P2P.\nThiết lập CSDL Neo4j và import dữ liệu chuẩn hóa.", "space_before": 4},
        {"text": "Tháng 06/2026:", "bold": True, "color": CYAN, "space_before": 14},
        {"text": "Phát triển backend FastAPI, GraphRAG và Circuit Breaker.\nXây dựng giao diện web ReactJS tương tác đồ thị.\nKiểm thử F1-Score hệ thống; viết báo cáo nghiệm thu.", "space_before": 4}
    ],
    top=Inches(1.2), left=Inches(6.8), width=Inches(5.9), height=Inches(5.4),
    glow_color=CYAN, emoji="📅", body_size=13.5
)


# ── SLIDE 11: 9. TÀI LIỆU THAM KHẢO (2-Column Grid) ────────────────
s = add_slide(prs)
add_top_bar(s)
add_bottom_bar(s, FOOTER)
add_title(s, "9. TÀI LIỆU THAM KHẢO")

# Column 1 (Ref 1 - 4)
draw_premium_card(
    s,
    "Tài liệu khoa học liên quan (1 - 4)",
    [
        {"text": "[1] H. Xu et al., \"medIKAL: Integrating Knowledge Graphs as Assistants of LLMs for Enhanced Clinical Diagnosis on EMRs,\" in Proceedings of the 31st International Conference on Computational Linguistics (COLING), Jan. 2025.", "size": 11.5, "space_before": 10},
        {"text": "[2] E. Evangelista et al., \"GraphRAG-Enabled Local Large Language Model for Gestational Diabetes Mellitus: Development of a Proof-of-Concept,\" JMIR Diabetes, vol. 11, e76454, Jan. 2026.", "size": 11.5, "space_before": 10},
        {"text": "[3] X. Wang et al., \"Building an intelligent diabetes Q&A system with knowledge graphs and large models,\" Frontiers in Public Health, vol. 13, 2025.", "size": 11.5, "space_before": 10},
        {"text": "[4] Y. Zheng et al., \"Automating Biomedical Knowledge Graph Construction For Context-Aware Scientific Inference,\" bioRxiv, Jan. 2026.", "size": 11.5, "space_before": 10}
    ],
    top=Inches(1.2), left=Inches(0.6), width=Inches(5.9), height=Inches(5.4),
    glow_color=CYAN, emoji="📚"
)

# Column 2 (Ref 5 - 8)
draw_premium_card(
    s,
    "Tài liệu khoa học liên quan (5 - 8)",
    [
        {"text": "[5] E. Zhang and H. Soh, \"Extract, Define, Canonicalize: An LLM-based Framework for Knowledge Graph Construction,\" in Proceedings of the 2024 Conference on Empirical Methods in Natural Language Processing, Miami, Florida, Nov. 2024, pp. 9548-9562.", "size": 11.5, "space_before": 10},
        {"text": "[6] J. Wu et al., \"Evidence-based Medical Large Language Model via Graph Retrieval-Augmented Generation,\" in Proceedings of the 63rd Annual Meeting of the Association for Computational Linguistics, 2025, pp. 28443-28467.", "size": 11.5, "space_before": 10},
        {"text": "[7] Y. Du, S. Li, A. Torralba, J. B. Tenenbaum, and I. Mordatch, \"Improving Factuality and Reasoning in Language Models through Multiagent Debate,\" in Proceedings of the 12th International Conference on Learning Representations (ICLR), May 2024.", "size": 11.5, "space_before": 10},
        {"text": "[8] C.-M. Chan et al., \"ChatEval: Towards better LLM-based evaluators through multi-agent debate,\" in Proceedings of the 12th International Conference on Learning Representations (ICLR), May 2024, Art. no. 412.", "size": 11.5, "space_before": 10}
    ],
    top=Inches(1.2), left=Inches(6.8), width=Inches(5.9), height=Inches(5.4),
    glow_color=EMERALD, emoji="📚"
)


# ── SLIDE 12: Thank You ───────────────────────────────────────────
s = add_slide(prs)
add_top_bar(s, Inches(0.12))
add_bottom_bar(s, "Trường Công Nghệ Số Và Trí Tuệ Nhân Tạo  ·  Trường Đại Học Nam Cần Thơ  ·  2026")

tb_thank = s.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.0))
tf_thank = tb_thank.text_frame
tf_thank.word_wrap = True

p_th1 = tf_thank.paragraphs[0]
p_th1.text = "XIN TRÂN TRỌNG CẢM ƠN HỘI ĐỒNG!"
p_th1.font.name = FONT_FAMILY
p_th1.font.size = Pt(38)
p_th1.font.bold = True
p_th1.font.color.rgb = WHITE
p_th1.alignment = PP_ALIGN.CENTER

p_th2 = tf_thank.add_paragraph()
p_th2.text = "Kính mong nhận được ý kiến đóng góp của quý Thầy/Cô\nđể đề cương đăng ký đề tài được hoàn thiện hơn."
p_th2.font.name = FONT_FAMILY
p_th2.font.size = Pt(20)
p_th2.font.color.rgb = GRAY_MUTED
p_th2.space_before = Pt(25)
p_th2.alignment = PP_ALIGN.CENTER

# ── Save ──────────────────────────────────────────────────────────
save_path = r"e:\HK2(2025-2026)\Đồ Án 2\MainFolder\MyProject\de_cuong_bao_cao_chuan.pptx"
prs.save(save_path)
print("OK - Ultra Premium Dark Presentation compiled successfully with Related Works Slide included!")
